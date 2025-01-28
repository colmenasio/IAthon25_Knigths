import os
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
import pdfkit

# Configuración de rutas
html_file = r"C:/Users/Celia/Desktop/IAthon/example.html"
docx_path = 'documento.docx'
pdf_path = 'documento.pdf'
ppt_path = 'documento.pptx'

# Leer el archivo HTML
try:
    with open(html_file, 'r', encoding='utf-8') as file:
        html_content = file.read()
except FileNotFoundError:
    print(f"Error: El archivo HTML '{html_file}' no existe.")
    exit()

# Procesar el HTML con BeautifulSoup
soup = BeautifulSoup(html_content, 'html.parser')

# Funciones auxiliares para Word
def apply_heading_style(paragraph, level):
    """Aplica estilos personalizados a encabezados."""
    run = paragraph.runs[0]
    if level == 1:
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0, 0, 0)  # Negro
        run.bold = True
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    elif level == 2:
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0, 51, 102)  # Azul oscuro
        run.bold = True
    elif level == 3:
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0, 102, 204)  # Azul claro
        run.bold = True

def apply_paragraph_style(paragraph):
    """Aplica estilos personalizados a los párrafos."""
    run = paragraph.runs[0]
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0, 0, 0)  # Negro
    paragraph.paragraph_format.space_after = Pt(12)  # Espaciado después del párrafo
    paragraph.paragraph_format.line_spacing = 1.5  # Espaciado entre líneas

# Crear documento Word
doc = Document()
for element in soup.find_all(['h1', 'h2', 'h3', 'p']):
    if element.name in ['h1', 'h2', 'h3']:
        level = int(element.name[1])
        paragraph = doc.add_heading(level=level)
        paragraph.add_run(element.text)
        apply_heading_style(paragraph, level)
    elif element.name == 'p':
        paragraph = doc.add_paragraph()
        paragraph.add_run(element.text)
        apply_paragraph_style(paragraph)

# Agregar imágenes locales al documento Word
for img_tag in soup.find_all('img'):
    img_src = img_tag.get('src')
    if img_src and not img_src.startswith('http'):
        try:
            doc.add_picture(img_src, width=Inches(4))  # Ajustar tamaño si es necesario
        except Exception as e:
            print(f"Error al agregar la imagen local '{img_src}': {e}")

# Guardar el documento Word
if os.path.exists(docx_path):
    os.remove(docx_path)
doc.save(docx_path)
print("Documento Word generado correctamente.")

# Conversión a PDF
try:
    config = pdfkit.configuration(wkhtmltopdf=r'C:/Program Files/wkhtmltopdf/bin/wkhtmltopdf.exe')
    if os.path.exists(pdf_path):
        os.remove(pdf_path)
    pdfkit.from_string(html_content, pdf_path, configuration=config)
    print("Documento PDF generado correctamente.")
except Exception as e:
    print(f"Error al generar el PDF: {e}")

# Crear presentación PowerPoint
prs = Presentation()
for element in soup.find_all(['h1', 'h2', 'h3', 'p']):
    if element.name == 'h1':
        slide_layout = prs.slide_layouts[0]  # Título
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = element.text
    elif element.name == 'h2':
        slide_layout = prs.slide_layouts[1]  # Título y contenido
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = element.text
    elif element.name in ['h3', 'p']:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        content = slide.shapes.placeholders[1]
        content.text += f"{element.text}\n"  # Agregar texto acumulativo

# Agregar imágenes locales a las diapositivas
for img_tag in soup.find_all('img'):
    img_src = img_tag.get('src')
    if img_src:
        try:
            slide = prs.slides[-1]
            slide.shapes.add_picture(img_src, Inches(1), Inches(1), width=Inches(4), height=Inches(3))
        except Exception as e:
            print(f"Error al agregar imagen en PowerPoint: {e}")

# Guardar presentación PowerPoint
if os.path.exists(ppt_path):
    os.remove(ppt_path)
prs.save(ppt_path)
print("Presentación PowerPoint generada correctamente.")
