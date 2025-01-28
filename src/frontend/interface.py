import streamlit as st
import pandas as pd
import PyPDF2
from pptx import Presentation
from bs4 import BeautifulSoup
import os
from docx import Document  # For Word document handling

# Initialize session state to track if analysis is done
if "analysis_done" not in st.session_state:
    st.session_state.analysis_done = False

logo_path = r"C:/Users/Celia/IAthon25_Knigths/src/frontend/logo_light.png"  # Path to the logo

# Custom CSS for styling inspired by ChatGPT interface
st.markdown("""
    <style>
        /* Main container */
        .main {
            background-color: #474444 !important;
            padding: 3rem !important;
            border-radius: 12px !important;
            box-shadow: 0 4px 10px rgba(0, 0, 0, 0.1) !important;
            font-family: 'Segoe UI', sans-serif !important;
            color: #474444 !important;
        }

        /* Title */
        h1 {
            color: #ec934f !important; 
            font-family: 'Segoe UI', sans-serif !important;
            font-size: 2.5rem !important;
            text-align: center !important; 
            margin-bottom: 1rem !important;
        }

        /* File uploader */
        .stFileUploader > div > div > div > button {
            background-color: #1e81b0 !important;
            color: white !important;
            border-radius: 8px !important;
            border: none !important;
            padding: 10px 20px !important;
            font-size: 1rem !important;
            font-family: 'Segoe UI', sans-serif !important;
            transition: background-color 0.3s ease !important;
        }

        .stFileUploader > div > div > div > button:hover {
            background-color: #1D4ED8 !important;
        }

        /* Select boxes */
        .stSelectbox > div > div > select {
            border-radius: 8px !important;
            border: 1px solid #E5E7EB !important;
            padding: 10px !important;
            font-size: 1rem !important;
            font-family: 'Segoe UI', sans-serif !important;
        }

        /* Buttons */
        .stButton > button {
            background-color: rgba(0,152,158,255) !important;
            color: white !important;
            border-radius: 8px !important;
            border: none !important;
            padding: 12px 24px !important;
            font-size: 1rem !important;
            font-family: 'Segoe UI', sans-serif !important;
            transition: background-color 0.3s ease !important;
        }

        .stButton > button:hover {
            background-color: rgba(100,152,158,255) !important;
        }

        /* Sidebar */
        .sidebar .sidebar-content {
            background-color: #1E293B !important;
            color: white !important;
            padding: 1rem !important;
            border-radius: 12px !important;
        }

        /* Markdown text */
        .stMarkdown {
            color: #0F172A !important;
            font-family: 'Segoe UI', sans-serif !important;
        }

        /* Error messages */
        .stAlert {
            color: white !important;
            border-radius: 8px !important;
            padding: 1rem !important;
        }

        /* Logo styling */
        .logo {
            top: 0px !important; /* Move the logo higher */
            left: 0px !important; /* Move the logo more to the left */
            width: 5px !important; /* Increase the size of the logo */
            height: auto !important; /* Maintain aspect ratio */
            z-index: 1000 !important; /* Ensure the logo stays on top */
        }
        
        .logo-container {
            position: absolute;
            top: 0px; /* Adjust this value to move higher */
            left: 0px; /* Adjust this value to move more to the left */
            z-index: 1000; /* Ensure it stays on top */
        }
        .logo-container img {
            width: 150px; /* Adjust the size of the logo */
            height: auto; /* Maintain aspect ratio */
        }
        <div class="logo-container">
            <img src="file://{logo_path}" alt="Logo">
        </div>
    """,
    unsafe_allow_html=True,
)

# Add logo to the upper left corner
logo_path = r"C:/Users/Celia/IAthon25_Knigths/src/frontend/logo_light.png"  # Relative path to the logo in the same folder
st.image(logo_path, use_container_width=True)  # Adjust width as needed

# Sidebar for additional options
with st.sidebar:
    st.markdown("<h2>Settings</h2>", unsafe_allow_html=True)
    language = st.selectbox("Select the language for analysis", ["English", "Spanish", "French", "German"])
    doc_type = st.selectbox("Select the document type", ["HTML", "PDF", "PowerPoint", "Word"])  # Added "Word" option

# Main content
st.markdown("### Upload your document", unsafe_allow_html=True)
uploaded_file = st.file_uploader("", type=['csv', 'xlsx', 'txt', 'html', 'pdf', 'pptx', 'docx'])  # Added 'docx'

# Text boxes for description and problem with placeholders
data_description = st.text_area(
    "Provide a brief description of the data:",
    placeholder="Enter a description of the data here...",
    value=""
)

problem_statement = st.text_area(
    "Describe the problem you want to solve with this data:",
    placeholder="Enter the problem statement here...",
    value=""
)

# Function to save the uploaded file locally
def save_uploaded_file(uploaded_file, save_directory):
    if not os.path.exists(save_directory):
        os.makedirs(save_directory)  # Create the directory if it doesn't exist
    file_path = os.path.join(save_directory, uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return file_path

# Function to save text to a file
def save_text_to_file(text, file_name, save_directory):
    if not os.path.exists(save_directory):
        os.makedirs(save_directory)  # Create the directory if it doesn't exist
    file_path = os.path.join(save_directory, file_name)
    with open(file_path, "w") as f:
        f.write(text)
    return file_path

# Function to analyze data
def analyze_data(data, language, doc_type):
    st.success(f"Analysis of {doc_type} document in {language} completed!")
    st.markdown("### Data View", unsafe_allow_html=True)
    if isinstance(data, pd.DataFrame):
        st.write(data.head())  # Display the first few rows for CSV/Excel files
    else:
        st.write(data)  # Display text content for other file types

# Button to perform analysis
if st.button('Analyze Data'):
    if uploaded_file is not None:
        # Save the uploaded file locally
        save_directory = "uploaded_files"  # Directory to save the uploaded files
        saved_file_path = save_uploaded_file(uploaded_file, save_directory)

        # Save the data description and problem statement to text files
        if data_description:
            description_file_path = save_text_to_file(data_description, "data_description.txt", save_directory)

        if problem_statement:
            problem_file_path = save_text_to_file(problem_statement, "problem_statement.txt", save_directory)

        # Process the file based on its type
        if uploaded_file.type == "text/csv":
            data = pd.read_csv(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            data = pd.read_excel(uploaded_file)
        elif uploaded_file.type == "text/plain":
            data = uploaded_file.read().decode("utf-8")
        elif uploaded_file.type == "text/html":
            soup = BeautifulSoup(uploaded_file, 'html.parser')
            data = soup.get_text()
        elif uploaded_file.type == "application/pdf":
            reader = PyPDF2.PdfFileReader(uploaded_file)
            data = ""
            for page in range(reader.numPages):
                data += reader.getPage(page).extract_text()
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(uploaded_file)
            data = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        data += shape.text + "\n"
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(uploaded_file)
            data = ""
            for para in doc.paragraphs:
                data += para.text + "\n"
        analyze_data(data, language, doc_type)

        # Set session state to indicate analysis is done
        st.session_state.analysis_done = True
    else:
        st.error("Please upload a file to analyze.")

# Add a button to download the generated document
if st.session_state.analysis_done:
    html_file = r"C:/Users/Celia/Desktop/IAthon/example.html"
    if os.path.exists(html_file):
        with open(html_file, "r", encoding="utf-8") as file:
            html_content = file.read()
        
        # Dynamically set the download button label based on the selected format
        download_label = f"Download Generated {doc_type}"
        st.download_button(
            label=download_label,
            data=html_content,
            file_name=f"generated_document.{doc_type.lower()}",  # Set file extension dynamically
            mime="text/html" if doc_type == "HTML" else 
                 "application/pdf" if doc_type == "PDF" else 
                 "application/vnd.openxmlformats-officedocument.presentationml.presentation" if doc_type == "PowerPoint" else 
                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document"  # For Word
        )
    else:
        st.warning("The HTML file does not exist at the specified path.")